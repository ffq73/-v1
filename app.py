%%writefile app.py
import streamlit as st
import docx
from pptx import Presentation
import re
import dashscope
from http import HTTPStatus

# --- 1. 基础工具函数 ---

def clean_text(text):
    """标准化清洗"""
    if not text: return ""
    return "".join(text.split())

def split_into_segments(full_text):
    """打散成短句集合"""
    segments = re.split(r'[。；！？\n]+', full_text)
    cleaned_segments = []
    for seg in segments:
        cleaned = clean_text(seg)
        if len(cleaned) > 2:
            cleaned_segments.append(cleaned)
    return set(cleaned_segments)

# --- 🟢 核心修复部分开始 🟢 ---

def get_docx_text(docx_file):
    """解析Word内容 (增强版：修复了合并单元格报错的问题)"""
    try:
        doc = docx.Document(docx_file)
    except Exception as e:
        st.error(f"Word文件损坏或无法读取: {e}")
        return set(), ""

    full_text = []
    
    # 1. 提取段落
    for para in doc.paragraphs:
        full_text.append(para.text)
        
    # 2. 提取表格 (增加容错机制)
    for i, table in enumerate(doc.tables):
        try:
            for row in table.rows:
                try:
                    # 尝试正常读取单元格
                    for cell in row.cells:
                        full_text.append(cell.text)
                except (ValueError, IndexError):
                    # 🚨 如果遇到合并单元格报错 (grid_offset error)
                    # 我们尝试直接读取该行 XML 中的文本，绕过网格计算
                    # 这是一个“暴力”读取法，能防止报错
                    try:
                        for cell in row._element.tc_lst:
                            for p in cell.p_lst:
                                # 简单拼接 xml 里的文本节点
                                t_nodes = p.xpath('.//w:t')
                                text_content = "".join([node.text for node in t_nodes if node.text])
                                full_text.append(text_content)
                    except:
                        # 如果还不行，就只能跳过这一行，保命要紧
                        pass
        except Exception:
            # 如果整个表格结构都坏了，跳过该表格
            continue

    merged_text = "\n".join(full_text)
    return split_into_segments(merged_text), merged_text

# --- 🔴 核心修复部分结束 🔴 ---

def get_pptx_text(pptx_file):
    """解析PPT内容"""
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        st.error(f"PPT文件损坏或无法读取: {e}")
        return set(), ""
        
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                full_text.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        full_text.append(cell.text)
    merged_text = "\n".join(full_text)
    return split_into_segments(merged_text), merged_text

# --- 2. AI 分析核心逻辑 ---

def call_qwen_analysis(api_key, base_context, suspicious_list):
    """调用通义千问 API"""
    dashscope.api_key = api_key
    
    items_str = "\n".join([f"- {item}" for item in suspicious_list])
    
    prompt = f"""
    你是一个严厉的行研报告审核员。
    【任务目标】
    判断【待审核内容】是否在【基准事实】中有依据。
    
    【基准事实 (Source)】:
    {base_context[:30000]} 
    
    【待审核内容 (Target - 疑似模板残留或错误)】:
    {items_str}

    【要求】
    1. 如果它是对基准事实的合理概括、同义改写，标记为【✅ 通过】。
    2. 如果基准事实中完全没有提及，标记为【❌ 存疑/残留】。
    请直接输出分析结果。
    """

    try:
        response = dashscope.Generation.call(
            model=dashscope.Generation.Models.qwen_turbo,
            prompt=prompt
        )
        if response.status_code == HTTPStatus.OK:
            return response.output.text
        else:
            return f"API 调用失败: {response.code} - {response.message}"
    except Exception as e:
        return f"发生错误: {str(e)}"

# --- 3. Streamlit 界面 ---

st.set_page_config(page_title="行研卫士 Pro", layout="wide")

st.sidebar.title("🤖 AI 设置")
api_key = st.sidebar.text_input("请输入 DashScope API Key", type="password")
st.sidebar.markdown("[👉 点击申请阿里云 Key](https://bailian.console.aliyun.com/)")

st.title("🛡️ 行研搬砖卫士 (Pro)")
st.caption("修复了复杂表格报错问题，支持 AI 语义分析")

col1, col2 = st.columns(2)
with col1:
    base_file = st.file_uploader("📂 1. 基准 Word", type=['docx'])
with col2:
    student_file = st.file_uploader("📂 2. 实习生 PPT", type=['pptx'])

if base_file and student_file:
    st.divider()
    
    with st.spinner("正在解析文件..."):
        base_set, base_raw = get_docx_text(base_file)
        student_set, student_raw = get_pptx_text(student_file)
        
        ghost_content = list(student_set - base_set)

    if not ghost_content:
        st.success("🎉 完美！PPT 内容与 Word 完全字符级匹配。")
    else:
        st.warning(f"⚠️ 发现 {len(ghost_content)} 处内容无法直接匹配。")
        
        if st.button("调用 AI 分析 (Qwen)", type="primary"):
            if not api_key:
                st.error("请先在左侧输入 API Key！")
            else:
                if len(ghost_content) > 50:
                    st.warning("差异项过多，仅分析前 50 条...")
                    ghost_content = ghost_content[:50]
                
                with st.spinner("AI 正在思考..."):
                    ai_result = call_qwen_analysis(api_key, base_raw, ghost_content)
                
                st.subheader("📋 分析结果")
                st.markdown(ai_result)

    with st.expander("🔍 原始差异列表"):
        st.write(ghost_content)